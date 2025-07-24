from pptx import Presentation

def parse_pptx(file_bytes):
    prs = Presentation(file_bytes)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    chunks = "\n".join(texts).split("\n\n")
    return [{"text": chunk.strip()} for chunk in chunks if chunk.strip()]
